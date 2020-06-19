# -*- coding: utf-8 -*-
"""
Created on Fri Apr 06 17:46:44 2018

@author: TAQDJO
"""

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import json
import utils as u
import input as i

c = i.shared['cl']
n = f'PPTPeriodicalTemplate {c}.pptx'
df = u.run_sql_query(u.full_path('Ann_report_sql','den_incidences_offspec.sql'), i.shared)

###########################

prs = Presentation(n)
slide = prs.slides[27]
shapes = slide.shapes
rows = len(df) + 1
cols = 10
left = Inches(0.7)
top = Inches(1.2)
width = Inches(2.0)
height = Inches(0.6)
table = shapes.add_table(rows, cols, left, top, width, height).table


#margin_top
#Top margin of cell.
#margin_bottom
#Bottom margin of cell.


# set column widths
table.columns[0].width = Inches(1.5)
table.columns[1].width = Inches(1.1)
table.columns[2].width = Inches(1.5)
table.columns[3].width = Inches(0.8)
table.columns[4].width = Inches(0.6)
table.columns[5].width = Inches(1.5)
table.columns[6].width = Inches(0.6)
table.columns[7].width = Inches(0.6)
table.columns[8].width = Inches(0.6)
table.columns[9].width = Inches(0.6)

# write column headings
table.cell(0, 0).text = 'Ship'
table.cell(0, 1).text = 'Job name'
table.cell(0, 2).text = 'Port'
table.cell(0, 3).text = 'Bunker date'
table.cell(0, 4).text = 'Ordered Grade'
table.cell(0, 5).text = 'Supplier'
table.cell(0, 5).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Center left text
table.cell(0, 6).text = 'Quantity'
table.cell(0, 6).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
table.cell(0, 7).text = 'BDN Density'
table.cell(0, 7).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
table.cell(0, 8).text = 'Tested Density'
table.cell(0, 8).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
table.cell(0, 9).text = 'Diff'
table.cell(0, 9).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
for i in range(10):
    table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(9)
    table.cell(0, i).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    table.cell(0, i).fill.solid()
    table.cell(0, i).fill.fore_color.rgb = RGBColor(59, 142, 222)
# write body cells
for i, j in df.iterrows():
    table.cell(i + 1, 0).text = j[0]
    table.cell(i + 1, 1).text = j[1]
    table.cell(i + 1, 2).text = j[2]
    table.cell(i + 1, 3).text = str(j[3])
    table.cell(i + 1, 4).text = j[4]
    table.cell(i + 1, 5).text = j[5]
    table.cell(i + 1, 5).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Center align text
    table.cell(i + 1, 6).text = j[6]
    table.cell(i + 1, 6).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
    table.cell(i + 1, 7).text = j[7]
    table.cell(i + 1, 7).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
    table.cell(i + 1, 8).text = j[8]
    table.cell(i + 1, 8).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
    table.cell(i + 1, 9).text = j[9]
    table.cell(i + 1, 9).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
for r1 in range(len(df) + 1):
    for r2 in range(10):
        table.cell(r1, r2).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(r1, r2).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        table.cell(r1, r2).margin_top = 0
        table.cell(r1, r2).margin_top = 0
        table.cell(r1, r2).margin_left = 20000
        # table.cell(r1, r2).fill.solid()
        # table.cell(r1, r2).fill.fore_color.rgb=RGBColor(59, 142, 222)

prs.save(n)
