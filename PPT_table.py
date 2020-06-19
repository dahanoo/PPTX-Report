# -*- coding: utf-8 -*-
"""
Created on Fri Apr 06 17:46:44 2018

@author: TAQDJO
"""

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import pandas as pd
import json
import utils as u
import input as i

c = i.shared['cl']
n = f'PPTPeriodicalTemplate {c}.pptx'

df = u.run_sql_query('Ann_report_sql/AlSiGT40Table.sql', i.shared)

df = df.dropna()
df['Al&Si (mg/kg)'] = df['Al&Si (mg/kg)'].astype(int)
df = df[df['Al&Si (mg/kg)'] > 40]
df['Al&Si (mg/kg)'] = df['Al&Si (mg/kg)'].astype(str)
df = df.reset_index(drop=True)


def low(l):
    return l.title()


df['Ship'] = df['Ship'].apply(low)
df['Supplier'] = df['Supplier'].apply(low)

columns = ["Port", "PCounts", "Supplier", "SCounts"]

ports = df['Port'].value_counts()
ports = ports.reset_index(drop=False)

suppliers = df['Supplier'].value_counts()
suppliers = suppliers.reset_index(drop=False)
combined = pd.merge(ports, suppliers, left_index=True, right_index=True)

df2 = df.groupby(['Port', 'Supplier'])['Ship'].count().reset_index()
df4 = df.groupby(['Port'])['Ship'].count().reset_index()
df4['PercS'] = (df4.Ship / df4.Ship.sum()) * 100.
df4['PercS'] = df4['PercS'].round(0)
df4['labels'] = df4.Port.str.cat(df4['PercS'].astype(str), sep=' ')

df3 = df2
df3 = df3.rename(columns={0: 'Port', 1: 'Supplier', 2: 'Counts'})
df3['PercS'] = (df3['Ship'] / df3['Ship'].sum()) * 100.
df3['PercS'] = df3['PercS'].round(2)
df3['labels'] = df3['Supplier'].str.cat(df3['PercS'].astype(str), sep=' ')

###########################

prs = Presentation(n)
slide = prs.slides[20]
shapes = slide.shapes
rows = len(df) + 1
cols = 8
left = Inches(0.7)
top = Inches(1.)
width = Inches(2.0)
height = Inches(0.6)
table = shapes.add_table(rows, cols, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(1.5)
table.columns[1].width = Inches(1.0)
table.columns[2].width = Inches(1.2)
table.columns[3].width = Inches(1.5)
table.columns[4].width = Inches(1.5)
table.columns[5].width = Inches(0.75)
table.columns[6].width = Inches(0.75)
table.columns[7].width = Inches(0.75)
# write column headings
table.cell(0, 0).text = 'Vessel'
table.cell(0, 1).text = 'IMO Number'
table.cell(0, 2).text = 'Bunker date'
table.cell(0, 3).text = 'Port'
table.cell(0, 4).text = 'Supplier'
table.cell(0, 5).text = 'Al\n(mg/kg)'
table.cell(0, 5).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
table.cell(0, 6).text = 'Si\n(mg/kg)'
table.cell(0, 6).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
table.cell(0, 7).text = 'Al & Si\n(mg/kg)'
table.cell(0, 7).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
for i in range(8):
    table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(9)
    table.cell(0, i).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    table.cell(0, i).fill.solid()
    table.cell(0, i).fill.fore_color.rgb = RGBColor(59, 142, 222)
# write body cells
for i, j in df.iterrows():
    table.cell(i + 1, 0).text = j[0]
    table.cell(i + 1, 1).text = j[1]
    table.cell(i + 1, 2).text = j[2]
    table.cell(i + 1, 3).text = j[3]
    table.cell(i + 1, 4).text = j[4]
    table.cell(i + 1, 5).text = j[5]
    table.cell(i + 1, 5).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
    table.cell(i + 1, 6).text = j[6]
    table.cell(i + 1, 6).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
    table.cell(i + 1, 7).text = j[7]
    table.cell(i + 1, 7).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align text
for r1 in range(len(df) + 1):
    for r2 in range(8):
        table.cell(r1, r2).text_frame.paragraphs[0].font.size = Pt(9)
        table.cell(r1, r2).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        table.cell(r1, r2).margin_top = 0
        table.cell(r1, r2).margin_top = 0
        table.cell(r1, r2).margin_left = 20000

prs.save(n)
