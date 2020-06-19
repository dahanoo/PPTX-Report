from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx import Presentation
import utils as u
import json
import input as i

c = i.shared['cl']
user = i.shared['u'].lower()
n = f'PPTPeriodicalTemplate {c}.pptx'
prs = Presentation(n)
slide = prs.slides[32]

presenters = {'sophia': {'name': 'Sophia Themelarou', 'pos': 'Senior Specialist', 'tel': "+30 21 9907 874",
                         'email': 'sophia.themelarou@lr.org'},
              'maria': {'name': 'Maria Kyratsoudi', 'pos': 'Senior Specialist', 'tel': "+30 211 9907 752",
                        'email': 'Maria.Kyratsoudi@lr.org'},
              'naeem': {'name': 'Naeem Javaid', 'pos': 'Senior Specialist', 'tel': "+44 330 414 0571",
                        'email': 'Naeem.Javaid@lr.org'},
              'das': {'name': 'Partha Das', 'pos': 'Senior Specialist', 'tel': "+65 3163 0658",
                      'email': 'Partha.Das@lr.org'},
              'andrew': {'name': 'Andrew Shaw', 'pos': 'Managing Director', 'tel': "+44 164 244 0991",
                         'email': 'Andrew.Shaw@lr.org'},
              'zenia': {'name': 'Zenia Said', 'pos': 'Senior Administrator', 'tel': "+30 211 9907 732",
                        'email': 'zenia.said@lr.org'}}

if user in presenters:
    a = presenters[user]['name']
    b = presenters[user]['pos']
    c = presenters[user]['tel']
    d = presenters[user]['email']
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"""Presentation by:
{a} - {b}
T : {c}
E : {d}


"""
else:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"""Presentation by:

T : 
E : 

"""
font = run.font
font.name = 'Source Sans Pro'
font.size = Pt(18)
font.bold = False
font.italic = None  # cause value to be inherited from theme
font.color.rgb = RGBColor(0xff, 0xff, 0xff)
prs.save(n)
