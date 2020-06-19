from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx import Presentation
import utils as u
import json
import input as i

c = i.shared['cl']
n = f'PPTPeriodicalTemplate {c}.pptx'
df = u.run_sql_query(u.full_path('Ann_report_sql','s_incidences_offspec.sql'), i.shared)
prs = Presentation(n)
slide = prs.slides[24]
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    text_frame = shape.text_frame

sl = df["Sulphur level"].str.contains("LOW")
sh = df["Sulphur level"].str.contains("HIGH")
if len(df) == 0:
    a = b = c = 'No samples have'
if len(sh) == 1:
    a = str(sh.sum()) + ' sample has'
if len(sh) > 1:
    a = str(sh.sum()) + ' samples have'
if len(sl) == 1:
    b = c = str(sl.sum()) + ' samples have'
if len(sl) > 1:
    b = c = str(sl.sum()) + ' samples have'

p = text_frame.paragraphs[0]
p.level = 1
run = p.add_run()
run.text = f"""{a} been noted to fail compliance with MARPOL Annex VI Reg. 14.1.2 (max. S content outside ECA- SOx at 3.50% m/m).

{b} been noted to fail compliance with MARPOL Annex VI, Reg. 14.4.3 and the EU Directive requirement for EU ports (max S content at 0.10% m/m).

{c} been noted to fail compliance with the EU Directive requirement for EU ports"""

font = run.font
font.name = 'Source Sans Pro'
font.size = Pt(18)
font.bold = False
font.italic = None  # cause value to be inherited from theme
font.color.rgb = RGBColor(0xff, 0xff, 0xff)
prs.save(n)
