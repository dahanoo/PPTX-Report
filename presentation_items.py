# -*- coding: utf-8 -*-
"""
Created on Fri Jan 18 11:00:36 2019

@author: Danny
"""
from pptx import Presentation
import input as i
c = i.shared['cl']
n = f'PPTPeriodicalTemplate {c}.pptx'
prs = Presentation(n)
slide = prs.slides[8]
for shape in slide.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))
for shape in slide.shapes:
    print('%s' % shape.shape_type)
    
    
    