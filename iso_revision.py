from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx import Presentation
import utils as u
import json
import input as i

c = i.shared['cl']
n = f'PPTPeriodicalTemplate {c}.pptx'
df = u.run_sql_query(u.full_path('Ann_report_sql', 'iso_revision.sql'), i.shared)
def slide_num(n, txt, pres):
    prs = Presentation(pres)
    slide = prs.slides[n]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
    a = df.iloc[0]['revision']
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = txt.format(str(a))
    font = run.font
    font.name = 'Source Sans Pro'
    font.size = Pt(16)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme
    font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    prs.save(pres)


slide_num(27, """The {}, maximum water content limit for residual fuels is 0.5% v/v.""", n)
slide_num(27, """\n\nIt is also stated in {} that, marine gas oils should contain no water.""", n)
slide_num(28, """{}, maximum TSP content is set at 0.10% m/m. Total sediment measurement aims at indicating the tendency of the fuel to sludge deposition.
""", n)
