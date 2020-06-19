# -*- coding: utf-8 -*-
"""
Created on Mon Feb 12 13:36:15 2018

@author: TAQDJO
"""
from pptx.util import Inches, Cm
from pptx import Presentation
import json
import input as i


def PPt_Presentation_Placeholder(oldFileName, newFileName, img, left, top,
                                 slide_no, width, height):
    prs = Presentation(oldFileName)
    slide = prs.slides[slide_no]
    pic = slide.shapes.add_picture(img, Inches(left), Inches(top),
                                   width=Inches(width), height=Inches(height))
    prs.save(newFileName)


def open_PowerPoint_Presentation(oldFileName, newFileName, img, left, top,
                                 slide_no, width=None, height=None):
    prs = Presentation(oldFileName)
    slide = prs.slides[slide_no]
    pic = slide.shapes.add_picture(img, left, top, width, height)
    prs.save(newFileName)


c = i.shared['cl']
o = 'PPTPeriodicalTemplateNew.pptx'
n = f'PPTPeriodicalTemplate {c}.pptx'

try:
    open_PowerPoint_Presentation(o, n, 'plots/sample_distribution.png',
                                 Inches(1.1), Inches(1.1), 3, None, None)
except FileNotFoundError:
    pass

try:
    open_PowerPoint_Presentation(n, n, 'plots/sample_statuses.png', Inches(1),
                                 Inches(1.1), 4, None, None)
except FileNotFoundError:
    pass

try:
    open_PowerPoint_Presentation(n, n, 'plots/plot status by ship.png',
                                 Inches(1), Inches(1), 5, None, None)
except FileNotFoundError:
    pass

try:
    PPt_Presentation_Placeholder(n, n, 'plots/gbl_clt_statuses.jpeg', 0.5,
                                 1.3, 6, 8, 4.5)
except FileNotFoundError:
    pass

try:
    open_PowerPoint_Presentation(n, n, 'plots/sample_status_reason_a.png',
                                 Inches(1), Inches(1.1), 7, None, None)
except FileNotFoundError:
    pass

try:
    open_PowerPoint_Presentation(n, n, 'plots/sample_status_reason_r.png',
                                 Inches(1), Inches(1.1), 8, None, None)
except FileNotFoundError:
    exec(compile(open('no_red_reports.py').read(), 'no_red_reports.py',
                 'exec'))

try:
    open_PowerPoint_Presentation(n, n, 'plots/Client_Supplier_amber.png',
                                 Inches(0.5), Inches(1.1), 9, None, None)
except FileNotFoundError:
    pass

try:
    open_PowerPoint_Presentation(n, n, 'plots/Client_Supplier_red.png',
                                 Inches(0.5), Inches(1.1), 10, None, None)
except FileNotFoundError:
    pass

try:
    open_PowerPoint_Presentation(n, n, 'plots/Port Visited.png', Inches(0.5),
                                 Inches(1.1), 11, None, None)
except FileNotFoundError:
    pass

try:
    open_PowerPoint_Presentation(n, n, 'plots/Supplier Visited.png',
                                 Inches(0.5), Inches(1.1), 12, None, None)
except FileNotFoundError:
    pass

try:
    PPt_Presentation_Placeholder(n, n, 'plots/Sankey_Ship_Supplier.png', 0.5,
                                 1, 13, 10, 4.5)
except FileNotFoundError:
    pass

try:
    PPt_Presentation_Placeholder(n, n, 'plots/AnalysisQuantityFO.png', 0.5,
                                 1.3, 14, 8.5, 4)
except FileNotFoundError:
    pass


PPt_Presentation_Placeholder(n, n, 'plots/AnalysisQuantityGO.png', 0.5, 1.3,
                             15, 8.5, 4)

PPt_Presentation_Placeholder(n, n, 'plots/Client_times.png', 0.5, 1.1, 16, 10,
                             4.5)

try:
    PPt_Presentation_Placeholder(n, n, 'plots/Global_client_times.png', 0.5,
                                 1.1, 17, 10, 4.5)
except FileNotFoundError:
    pass

PPt_Presentation_Placeholder(n, n, 'plots/PPTOffSpecPerc1.png', 1, 1, 18,
                             4.2, 4.7)
PPt_Presentation_Placeholder(n, n, 'plots/PPTOffSpecPerc2.png', 6.1, 1, 18,
                             4.2, 4.7)

try:
    exec(compile(open('PPT_table.py').read(), 'PPT_table.py', 'exec'))
except AttributeError:
    pass

try:
    PPt_Presentation_Placeholder(n, n, 'plots/PPTPortSupplier.png', 0.5, 0.8,
                                 21, 8.5, 5)
except FileNotFoundError:
    pass

try:
    PPt_Presentation_Placeholder(n, n, 'plots/ppt_spider.jpg', 3.1, 1.2, 31,
                                 5, 4)
except FileNotFoundError:
    pass
try:
    PPt_Presentation_Placeholder(n, n, 'plots/AlsiDistPPT.png', 0.5, 1.4, 22,
                                 5.18, 3.59)
except FileNotFoundError:
    pass

try:
    PPt_Presentation_Placeholder(n, n, 'plots/dist_plot.jpg', 6.3, 1.4, 22,
                                 4.06, 3.83)
except FileNotFoundError:
    pass
