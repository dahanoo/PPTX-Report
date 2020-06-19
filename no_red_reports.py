from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx import Presentation
import json
import input as i

c = i.shared['cl']
n = f'PPTPeriodicalTemplate {c}.pptx'


def slide_num(sl, txt, pres):
    prs = Presentation(pres)
    slide = prs.slides[sl]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = txt
    font = run.font
    font.name = 'Source Sans Pro'
    font.size = Pt(44)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme
    font.color.rgb = RGBColor(59, 142, 222)
    prs.save(pres)


slide_num(8, """No red reports""", n)
slide_num(10, """No red reports""", n)
